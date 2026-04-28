"""
Board presentation — modern dark theme.
Run: python scenario-04/generate_ppt.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
BG          = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
SURFACE     = RGBColor(0x16, 0x2B, 0x40)   # card background
SURFACE2    = RGBColor(0x1C, 0x35, 0x4D)   # elevated card
CYAN        = RGBColor(0x00, 0xC2, 0xCB)   # primary accent
CYAN_DIM    = RGBColor(0x00, 0x7A, 0x82)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY1       = RGBColor(0xC8, 0xD8, 0xE8)   # body text
GRAY2       = RGBColor(0x6E, 0x8C, 0xA8)   # muted
RED         = RGBColor(0xFF, 0x4D, 0x6D)
ORANGE      = RGBColor(0xFF, 0x9F, 0x1C)
GREEN       = RGBColor(0x06, 0xD6, 0xA0)
PURPLE      = RGBColor(0x9B, 0x5D, 0xE5)

W = Inches(13.33)
H = Inches(7.5)

ROUNDED = 5   # MSO autoshape id for rounded rectangle


# ── Primitives ────────────────────────────────────────────────────────────────
def prs_new() -> Presentation:
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def fill_bg(slide, color=BG):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, rounded=False, line=None):
    shape_id = ROUNDED if rounded else 1
    s = slide.shapes.add_shape(shape_id, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    if rounded:
        s.adjustments[0] = 0.04
    return s

def add_text(slide, l, t, w, h, text, size=14, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def tag(slide, l, t, text, color=CYAN):
    """Small uppercase label tag."""
    add_text(slide, l, t, Inches(5), Inches(0.32), text,
             size=9, bold=True, color=color)

def divider(slide, t, color=SURFACE2):
    add_rect(slide, Inches(0.55), t, Inches(12.23), Pt(1.5), color)

def accent_bar(slide, l, t, h=Inches(0.05), color=CYAN):
    add_rect(slide, l, t, Inches(12.23), h, color)

# ── Card helpers ──────────────────────────────────────────────────────────────
def card(slide, l, t, w, h, top_color=None):
    add_rect(slide, l, t, w, h, SURFACE, rounded=True)
    if top_color:
        add_rect(slide, l, t, w, Inches(0.055), top_color, rounded=False)

def kpi_card(slide, l, t, w, h, label, value, val_color=CYAN, sub=None):
    card(slide, l, t, w, h, top_color=val_color)
    add_text(slide, l+Inches(0.2), t+Inches(0.15), w-Inches(0.4), Inches(0.3),
             label.upper(), size=8, bold=True, color=GRAY2)
    add_text(slide, l+Inches(0.2), t+Inches(0.45), w-Inches(0.4), Inches(0.65),
             value, size=26, bold=True, color=val_color)
    if sub:
        add_text(slide, l+Inches(0.2), t+Inches(1.05), w-Inches(0.4), Inches(0.28),
                 sub, size=9, color=GRAY2, italic=True)

def section_header(slide, label, title, subtitle=None):
    """Top band with section label + title."""
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.25), SURFACE)
    accent_bar(slide, Inches(0.55), Inches(1.23))
    tag(slide, Inches(0.55), Inches(0.18), label)
    add_text(slide, Inches(0.55), Inches(0.48), Inches(12), Inches(0.65),
             title, size=22, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.55), Inches(0.95), Inches(12), Inches(0.35),
                 subtitle, size=11, color=GRAY2)

def numbered_row(slide, t, num, title, desc, timing=None, accent_color=CYAN):
    """One step row for next-steps slide."""
    # number circle
    add_rect(slide, Inches(0.55), t+Inches(0.08), Inches(0.52), Inches(0.52),
             accent_color, rounded=True)
    add_text(slide, Inches(0.55), t+Inches(0.1), Inches(0.52), Inches(0.5),
             str(num), size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.25), t+Inches(0.04), Inches(7.5), Inches(0.38),
             title, size=13, bold=True, color=WHITE)
    add_text(slide, Inches(1.25), t+Inches(0.4), Inches(9.8), Inches(0.35),
             desc, size=10, color=GRAY1)
    if timing:
        add_rect(slide, Inches(11.7), t+Inches(0.12), Inches(1.3), Inches(0.38),
                 accent_color, rounded=True)
        add_text(slide, Inches(11.7), t+Inches(0.14), Inches(1.3), Inches(0.34),
                 timing, size=9, bold=True, color=BG, align=PP_ALIGN.CENTER)


# ── Slide 01 — Cover ──────────────────────────────────────────────────────────
def s01_cover(prs):
    s = blank(prs)
    fill_bg(s)

    # left accent strip
    add_rect(s, Inches(0), Inches(0), Inches(0.22), H, CYAN)

    # big title block
    add_text(s, Inches(0.65), Inches(1.4), Inches(8), Inches(0.5),
             "CHURN RATE ANALYTICS", size=11, bold=True, color=CYAN)
    add_text(s, Inches(0.65), Inches(1.9), Inches(9), Inches(2.0),
             "One Metric.\nOne Definition.\nOne Dashboard.", size=40,
             bold=True, color=WHITE)

    divider(s, Inches(4.15), color=CYAN_DIM)

    add_text(s, Inches(0.65), Inches(4.4), Inches(8), Inches(0.45),
             "Sostituisce 40 dashboard su 3 BI tool con una sola fonte di verità.",
             size=14, color=GRAY1)
    add_text(s, Inches(0.65), Inches(5.0), Inches(8), Inches(0.38),
             "Definizione canonica v1.0 — approvata 28 aprile 2026",
             size=11, color=GRAY2)

    # right stat block
    for i, (val, lbl) in enumerate([
        ("40", "dashboard eliminati"),
        ("4",  "stakeholder allineati"),
        ("1",  "fonte di verità"),
    ]):
        t = Inches(2.2) + i * Inches(1.45)
        kpi_card(s, Inches(10.2), t, Inches(2.7), Inches(1.2), lbl, val, CYAN)

    add_text(s, Inches(0.65), Inches(6.85), Inches(5), Inches(0.38),
             "Scenario 04  ·  Hackathon", size=9, color=GRAY2)


# ── Slide 02 — Il problema ────────────────────────────────────────────────────
def s02_problem(prs):
    s = blank(prs)
    fill_bg(s)
    section_header(s, "01 — IL PROBLEMA",
                   "Stessa domanda, 4 risposte diverse")

    q = "« Qual è il churn rate di Q2 2025? »"
    add_text(s, Inches(0.55), Inches(1.45), Inches(12.2), Inches(0.5),
             q, size=16, italic=True, color=CYAN)

    answers = [
        ("VP Sales",   "3.1%",  "Logo churn\nsenza downgrade",        ORANGE),
        ("CS Head",    "2.4%",  "CS-accountable\nescluse bankruptcies", PURPLE),
        ("Finance",    "5.8%",  "nMRR grezzo\n+ 340 duplicati CRM",    RED),
        ("Analyst",    "2.1%",  "nMRR canonical D4\nclean",            GREEN),
    ]
    cw = Inches(2.85)
    for i, (role, val, method, col) in enumerate(answers):
        l = Inches(0.55) + i * (cw + Inches(0.22))
        card(s, l, Inches(2.1), cw, Inches(3.5), top_color=col)
        add_text(s, l+Inches(0.2), Inches(2.3), cw-Inches(0.4), Inches(0.38),
                 role.upper(), size=9, bold=True, color=col)
        add_text(s, l+Inches(0.2), Inches(2.7), cw-Inches(0.4), Inches(0.9),
                 val, size=44, bold=True, color=col)
        add_text(s, l+Inches(0.2), Inches(3.65), cw-Inches(0.4), Inches(0.7),
                 method, size=11, color=GRAY1)

    # warning banner
    add_rect(s, Inches(0.55), Inches(5.85), Inches(12.23), Inches(1.3),
             RGBColor(0x2A, 0x10, 0x15), rounded=True)
    add_rect(s, Inches(0.55), Inches(5.85), Inches(0.06), Inches(1.3), RED)
    add_text(s, Inches(0.8), Inches(5.95), Inches(11.7), Inches(0.38),
             "BUG CONFERMATO — Q2 2025", size=11, bold=True, color=RED)
    add_text(s, Inches(0.8), Inches(6.33), Inches(11.7), Inches(0.65),
             "Il 5.8% presentato al board era causato da 340 eventi duplicati (CRM retry storm). "
             "Valore corretto: 2.1%.  Delta: +3.7pp falsi.  Rettifica formale ancora da emettere.",
             size=10, color=GRAY1)


# ── Slide 03 — Stakeholder requirements ──────────────────────────────────────
def s03_stakeholders(prs):
    s = blank(prs)
    fill_bg(s)
    section_header(s, "02 — REQUISITI",
                   "Cosa voleva ogni stakeholder",
                   "12 punti di disaccordo → 1 definizione approvata da tutti")

    rows = [
        ("VP Sales & Marketing", ORANGE,
         "Vista logocentrica",
         "Churn = clienti persi. Toggle CS-accountable separato dal KPI finanziario."),
        ("CS Head", PURPLE,
         "Esclusioni strutturali",
         "Bankruptcies e M&A fuori dal suo KPI. Save reversal entro 30gg = churn annullato."),
        ("Finance Director", GREEN,
         "nMRR normalizzato",
         "Annuali ÷ 12. Micro-contratti < €200/m esclusi dall'headline. Versioning obbligatorio."),
        ("Data Analyst", CYAN,
         "Audit trail completo",
         "Definizione con soglie numeriche esplicite, boundary cases, test automatici, export CSV."),
    ]

    rh = Inches(1.2)
    for i, (name, col, badge, desc) in enumerate(rows):
        t = Inches(1.5) + i * (rh + Inches(0.06))
        card(s, Inches(0.55), t, Inches(12.23), rh, top_color=col)
        add_rect(s, Inches(0.55), t, Inches(0.06), rh, col)

        add_rect(s, Inches(0.85), t+Inches(0.35), Inches(1.8), Inches(0.38),
                 col, rounded=True)
        add_text(s, Inches(0.85), t+Inches(0.37), Inches(1.8), Inches(0.34),
                 badge, size=9, bold=True, color=BG, align=PP_ALIGN.CENTER)

        add_text(s, Inches(2.85), t+Inches(0.08), Inches(3.5), Inches(0.38),
                 name, size=13, bold=True, color=WHITE)
        add_text(s, Inches(2.85), t+Inches(0.52), Inches(9.6), Inches(0.5),
                 desc, size=11, color=GRAY1)


# ── Slide 04 — Definizione v1.0 ───────────────────────────────────────────────
def s04_definition(prs):
    s = blank(prs)
    fill_bg(s)
    section_header(s, "03 — DEFINIZIONE CANONICA v1.0",
                   "Ogni soglia è numerica. Nessun termine vago.",
                   "Approvata da VP · CS · Finance · Analyst  —  28 apr 2026")

    rules = [
        ("Formula",           "nMRR Lost  ÷  nMRR Active (inizio periodo)  × 100",          CYAN),
        ("Data canonica",     "contract_end_date  — non la data di cancellazione",           CYAN),
        ("nMRR",              "Annuali ÷ 12  —  nessun picco artificiale a scadenza",        CYAN),
        ("Micro-contratti",   "< €200/mese → esclusi dall'headline · tracciati separatamente", ORANGE),
        ("Downgrade",         "Conta se delta ≥ €50  E  ≥ 10% del nMRR precedente",          ORANGE),
        ("Save reversal",     "Riattivazione entro 30 giorni da contract_end_date → annullato", GREEN),
        ("Grace post-rinnovo","Cancellazione entro 14gg da rinnovo automatico → rimborso, non churn", GREEN),
        ("Pausa",             "At-risk ≤ 90 giorni · oltre = churn automatico (pause_expiry)", RED),
    ]

    col_w = Inches(5.9)
    for i, (label, rule, col) in enumerate(rules):
        ci = i % 2
        ri = i // 2
        l = Inches(0.55) + ci * (col_w + Inches(0.22))
        t = Inches(1.55) + ri * Inches(1.32)
        card(s, l, t, col_w, Inches(1.2))
        add_rect(s, l, t, col_w, Inches(0.05), col)
        add_text(s, l+Inches(0.2), t+Inches(0.1), col_w-Inches(0.4), Inches(0.3),
                 label.upper(), size=8, bold=True, color=col)
        add_text(s, l+Inches(0.2), t+Inches(0.42), col_w-Inches(0.4), Inches(0.65),
                 rule, size=11, color=GRAY1)


# ── Slide 05 — Architettura ───────────────────────────────────────────────────
def s05_architecture(prs):
    s = blank(prs)
    fill_bg(s)
    section_header(s, "04 — ARCHITETTURA",
                   "Tre layer indipendenti — testati, versionati, auditabili")

    layers = [
        ("DATA",          ORANGE, "5 CSV con noise iniettato",
         ["340 duplicati CRM", "Timezone mismatch UTC/CET",
          "15 segmenti errati", "12 end_date mancanti"]),
        ("ENGINE",        CYAN,   "FastAPI · Pydantic · Python",
         ["calculator.py — pure functions", "18 unit test · 100% pass",
          "Ogni risultato tagged definition_version", "_window_bounds · save reversals · micro filter"]),
        ("DASHBOARD",     GREEN,  "Streamlit · Plotly · Claude",
         ["KPI cards · trend · breakdown", "At-risk con risk score e CSM",
          "NL eval harness — 20/20 pass", "Export CSV versionato"]),
    ]

    lw = Inches(3.9)
    for i, (name, col, sub, bullets) in enumerate(layers):
        l = Inches(0.55) + i * (lw + Inches(0.22))
        card(s, l, Inches(1.5), lw, Inches(4.7), top_color=col)
        add_text(s, l+Inches(0.22), Inches(1.65), lw-Inches(0.44), Inches(0.5),
                 name, size=18, bold=True, color=col)
        add_text(s, l+Inches(0.22), Inches(2.1), lw-Inches(0.44), Inches(0.35),
                 sub, size=10, color=GRAY2)
        divider(s, Inches(2.5), SURFACE2)
        for j, b in enumerate(bullets):
            add_text(s, l+Inches(0.22), Inches(2.65)+j*Inches(0.62),
                     lw-Inches(0.44), Inches(0.5),
                     "· " + b, size=10, color=GRAY1)

        # arrow between layers
        if i < 2:
            ax = l + lw + Inches(0.04)
            add_text(s, ax, Inches(3.6), Inches(0.18), Inches(0.38),
                     "→", size=18, bold=True, color=GRAY2, align=PP_ALIGN.CENTER)

    # bottom KPI strip
    metrics = [("18/18", "Unit test"), ("20/20", "Golden Qs"),
               ("100%", "Accuracy"), ("0%", "False confidence")]
    mw = Inches(2.8)
    for i, (val, lbl) in enumerate(metrics):
        l = Inches(0.55) + i * (mw + Inches(0.2))
        kpi_card(s, l, Inches(6.35), mw, Inches(0.9), lbl, val, GREEN)


# ── Slide 06 — Riconciliazione storica ────────────────────────────────────────
def s06_reconciliation(prs):
    s = blank(prs)
    fill_bg(s)
    section_header(s, "05 — RICONCILIAZIONE STORICA",
                   "5 anni · 4 definizioni · 1 verità canonica",
                   "Ogni trimestre storico ricalcolato con D4 — le colonne precedenti non sono comparabili")

    defs = [
        ("D0", "2021 – Q1 2023",  "Logo churn puro",           GRAY2,
         "Misura solo i clienti persi. Non cattura la perdita di fatturato."),
        ("D1", "Q2 – Q4 2023",    "Logo + downgrade",          ORANGE,
         "Cambiato con il nuovo CFO. Include contrazioni ma senza soglia chiara."),
        ("D2", "2024 – Q1 2025",  "nMRR grezzo",               CYAN,
         "Migliora la misura revenue ma manca la normalizzazione annuali÷12."),
        ("D3", "Q2 2025",         "nMRR + duplicati CRM",      RED,
         "BUG: 340 eventi duplicati portano il rate da 2.1% a 5.8% (+3.7pp falsi)."),
        ("D4", "v1.0 canonica",   "nMRR normalizzato · clean", GREEN,
         "Definizione approvata. Retroattiva su 5 anni. Questo è il numero corretto."),
    ]

    rh = Inches(0.98)
    for i, (code, period, name, col, note) in enumerate(defs):
        t = Inches(1.5) + i * (rh + Inches(0.05))
        is_last = (i == len(defs) - 1)
        bg_col = SURFACE2 if is_last else SURFACE
        card(s, Inches(0.55), t, Inches(12.23), rh, top_color=col)

        # code badge
        add_rect(s, Inches(0.55), t, Inches(0.06), rh, col)
        add_rect(s, Inches(0.75), t+Inches(0.25), Inches(0.6), Inches(0.42),
                 col, rounded=True)
        add_text(s, Inches(0.75), t+Inches(0.27), Inches(0.6), Inches(0.38),
                 code, size=11, bold=True, color=BG, align=PP_ALIGN.CENTER)

        add_text(s, Inches(1.55), t+Inches(0.08), Inches(2.0), Inches(0.38),
                 period, size=10, color=GRAY2)
        add_text(s, Inches(3.6), t+Inches(0.08), Inches(3.2), Inches(0.38),
                 name, size=12, bold=True, color=WHITE if not is_last else col)
        add_text(s, Inches(6.9), t+Inches(0.08), Inches(5.7), Inches(0.75),
                 note, size=10, color=GRAY1 if not is_last else GREEN)

        if is_last:
            add_rect(s, Inches(12.4), t+Inches(0.25), Inches(0.38), Inches(0.42),
                     GREEN, rounded=True)
            add_text(s, Inches(12.4), t+Inches(0.27), Inches(0.38), Inches(0.38),
                     "✓", size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)


# ── Slide 07 — Dashboard ──────────────────────────────────────────────────────
def s07_dashboard(prs):
    s = blank(prs)
    fill_bg(s)
    section_header(s, "06 — LA DASHBOARD",
                   "Single source of truth — sostituisce 40 dashboard su 3 BI tool")

    kpis = [
        ("nMRR Churn Rate", "0.19%", CYAN),
        ("MRR Lost",        "€2.446",  CYAN),
        ("MRR Active",      "€1.3M",   GRAY2),
        ("Logo Churn",      "0.31%",   CYAN),
        ("CS Saves",        "0",       GRAY2),
        ("nMRR rischio 🔴", "€0",      GREEN),
    ]
    kw = Inches(2.0)
    for i, (lbl, val, col) in enumerate(kpis):
        kpi_card(s, Inches(0.55)+i*(kw+Inches(0.1)), Inches(1.5),
                 kw, Inches(1.15), lbl, val, col)

    features = [
        (CYAN,   "Vista Financial / Contractual", "Toggle nMRR ↔ logo churn"),
        (ORANGE, "Window: Monthly · R30 · R90",   "Flessibilità temporale on-the-fly"),
        (PURPLE, "Filtro segmento",                "Enterprise · Mid-market · SMB"),
        (CYAN,   "Toggle CS-accountable",          "Separa KPI CS da quello finanziario"),
        (GREEN,  "Contratti a rischio",             "Scadenza ≤90gg · risk score · CSM owner"),
        (GREEN,  "Export CSV versionato",           "Tutti i periodi con definition_version"),
        (ORANGE, "Glossario inline",                "Definizione v1.0 sempre visibile"),
        (PURPLE, "Storico 5 anni",                  "Confronto D0→D4 con annotazioni"),
    ]

    cw = Inches(5.9)
    for i, (col, feat, desc) in enumerate(features):
        ci, ri = i % 2, i // 2
        l = Inches(0.55) + ci * (cw + Inches(0.22))
        t = Inches(2.85) + ri * Inches(0.96)
        card(s, l, t, cw, Inches(0.85))
        add_rect(s, l, t+Inches(0.18), Inches(0.05), Inches(0.5), col)
        add_text(s, l+Inches(0.22), t+Inches(0.1), Inches(2.5), Inches(0.35),
                 feat, size=11, bold=True, color=WHITE)
        add_text(s, l+Inches(0.22), t+Inches(0.48), Inches(5.4), Inches(0.3),
                 desc, size=10, color=GRAY2)


# ── Slide 08 — Eval harness ───────────────────────────────────────────────────
def s08_eval(prs):
    s = blank(prs)
    fill_bg(s)
    section_header(s, "07 — NL EVAL HARNESS",
                   "Claude risponde a domande sul churn con tool use — 20 golden questions")

    results = [
        ("Accuracy",         "100%", "> 80%",  GREEN),
        ("Refusal accuracy", "100%", "> 90%",  GREEN),
        ("False confidence", "0%",   "< 5%",   GREEN),
        ("Pass rate",        "100%", "> 80%",  GREEN),
    ]
    rw = Inches(2.85)
    for i, (lbl, val, tgt, col) in enumerate(results):
        kpi_card(s, Inches(0.55)+i*(rw+Inches(0.17)), Inches(1.5),
                 rw, Inches(1.4), lbl, val, col, sub=f"Target {tgt}")

    add_text(s, Inches(0.55), Inches(3.1), Inches(6), Inches(0.4),
             "DISTRIBUZIONE PER TIPO", size=9, bold=True, color=GRAY2)

    types = [
        ("answerable",  4, CYAN,   "Lookup metriche — get_metric"),
        ("comparison",  3, ORANGE, "Confronto periodi e segmenti"),
        ("definition",  5, PURPLE, "Regole e soglie della definizione"),
        ("edge_case",   3, GREEN,  "Boundary cases (grace, micro, save)"),
        ("refusal",     5, RED,    "Rifiuto corretto (forecast, PII, OOB)"),
    ]

    max_count = 5
    bw_unit = Inches(1.6)
    for i, (qtype, count, col, desc) in enumerate(types):
        t = Inches(3.6) + i * Inches(0.68)
        bw = bw_unit * count
        add_rect(s, Inches(0.55), t, bw, Inches(0.48), col, rounded=True)
        add_text(s, Inches(0.55)+bw+Inches(0.15), t+Inches(0.06),
                 Inches(8), Inches(0.38),
                 f"{qtype}  ({count})  —  {desc}", size=11, color=GRAY1)

    # tools
    add_rect(s, Inches(7.0), Inches(3.1), Inches(5.8), Inches(3.85),
             SURFACE, rounded=True)
    add_rect(s, Inches(7.0), Inches(3.1), Inches(5.8), Inches(0.05), CYAN)
    add_text(s, Inches(7.2), Inches(3.22), Inches(5.4), Inches(0.35),
             "TOOL USE", size=9, bold=True, color=CYAN)
    tools = ["get_metric", "compare_periods", "list_definitions", "explain_calculation"]
    for i, t_name in enumerate(tools):
        t = Inches(3.7) + i * Inches(0.72)
        add_rect(s, Inches(7.2), t, Inches(2.2), Inches(0.42), SURFACE2, rounded=True)
        add_text(s, Inches(7.35), t+Inches(0.06), Inches(2.0), Inches(0.35),
                 t_name, size=10, bold=True, color=CYAN)

    add_text(s, Inches(7.2), Inches(6.1), Inches(5.4), Inches(0.35),
             "Prompt caching · agentic loop ≤5 turns", size=9, color=GRAY2)


# ── Slide 09 — Next steps ─────────────────────────────────────────────────────
def s09_next(prs):
    s = blank(prs)
    fill_bg(s)
    section_header(s, "08 — NEXT STEPS",
                   "Dalla demo alla produzione")

    steps = [
        (RED,    "Rettifica formale Q2 2025",
         "Emettere la nota al board: churn reale 2.1%, non 5.8%. Delta: +3.7pp. Aggiornare i verbali.",
         "Immediato"),
        (ORANGE, "Deploy su server interno",
         "Docker container su infrastruttura aziendale. Accesso via browser, nessuna installazione.",
         "2 settimane"),
        (CYAN,   "Connessione DB di produzione",
         "Sostituire i CSV con lettura dal warehouse. data_loader.py resta l'unico punto di accesso.",
         "1 mese"),
        (GREEN,  "Deprecazione dashboard legacy",
         "30 giorni di run parallelo, poi dismissione dei 40 dashboard. Comunicare il nuovo URL unico.",
         "3 mesi"),
        (PURPLE, "Definizione v1.1",
         "Revisione annuale. Aggiornare DEFINITION_VERSION — footer e glossario si aggiornano in automatico.",
         "Apr 2027"),
    ]

    for i, (col, title, desc, timing) in enumerate(steps):
        numbered_row(s, Inches(1.4)+i*Inches(1.12), i+1, title, desc, timing, col)


# ── Main ──────────────────────────────────────────────────────────────────────
def main():
    prs = prs_new()
    s01_cover(prs)
    s02_problem(prs)
    s03_stakeholders(prs)
    s04_definition(prs)
    s05_architecture(prs)
    s06_reconciliation(prs)
    s07_dashboard(prs)
    s08_eval(prs)
    s09_next(prs)

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "churn_rate_board_presentation.pptx")
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
